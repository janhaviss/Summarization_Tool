import logging
from typing import Union
from pathlib import Path
import tempfile
import asyncio
from fastapi import UploadFile, HTTPException, status
from schemas.summarization import SummaryRequest
import PyPDF2
from pptx import Presentation
import docx
from pydantic_settings import BaseSettings
import re
from service.config import settings

# Set up logging configuration
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)



ALLOWED_FILE_TYPES = {
    'application/pdf': '.pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx'
}

class SummarizationService:
    def __init__(self):
        self.summarizer_pipeline = None
        self.model_loaded = False
        self.model_loading_lock = asyncio.Lock()

    async def initialize(self):
        """Lazy loading of models with thread safety"""
        if not self.model_loaded:
            async with self.model_loading_lock:
                if not self.model_loaded:
                    await self._load_models()
                    self.model_loaded = True

    async def _load_models(self):
        """Load summarization models with better error handling"""
        try:
            from transformers import pipeline
            logger.info("Loading summarization model...")
            self.summarizer_pipeline = pipeline(
                "summarization",
                model=settings.MODEL_NAME,
                device="cpu"
            )
            logger.info("Model loaded successfully")
        except ImportError as e:
            logger.error(f"Failed to load transformers: {str(e)}")
            raise RuntimeError(f"Failed to load transformers: {str(e)}")
        except Exception as e:
            logger.error(f"Model loading failed: {str(e)}")
            raise RuntimeError(f"Model loading failed: {str(e)}")

    async def summarize_content(
        self,
        content: Union[SummaryRequest, UploadFile],
        method: str = "transformers",
        **kwargs
    ) -> str:
        await self.initialize()

        if self.summarizer_pipeline is None:
            raise HTTPException(
                status_code=500,
                detail="Summarization model not initialized"
            )

        if isinstance(content, SummaryRequest):
            return await self.summarize_text(content.text, method, **kwargs)

        elif isinstance(content, UploadFile):
            return await self.process_uploaded_file(content, method, **kwargs)

        else:
            raise HTTPException(
                status_code=400,
                detail="Invalid input type - must be text or file"
            )

    async def process_uploaded_file(
        self,
        file: UploadFile,
        method: str,
        **kwargs
    ) -> str:
        file_ext = Path(file.filename).suffix.lower()
        temp_file_path = None

        try:
            contents = await file.read()

            if len(contents) > settings.max_file_size_mb * 1024 * 1024:
                raise HTTPException(
                    status_code=413,
                    detail=f"File too large (>{settings.max_file_size_mb}MB)"
                )

            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                temp_file.write(contents)
                temp_file_path = temp_file.name

            if file_ext == '.pdf':
                text = await self.extract_text_from_pdf(temp_file_path)
            elif file_ext == '.docx':
                text = await self.extract_text_from_docx(temp_file_path)
            elif file_ext == '.pptx':
                text = await self.extract_text_from_pptx(temp_file_path)
            else:
                raise HTTPException(
                    status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE,
                    detail="Unsupported file type"
                )

            text = self.clean_text(text)
            if len(text.split()) > settings.max_text_length:
                text = " ".join(text.split()[:settings.max_text_length])

            return await self.summarize_text(text, method, **kwargs)

        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"File processing error: {str(e)}"
            )
        finally:
            if temp_file_path:
                Path(temp_file_path).unlink(missing_ok=True)

    async def extract_text_from_pdf(self, file_path: str) -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._sync_extract_pdf_text, file_path)

    def _sync_extract_pdf_text(self, file_path: str) -> str:
        text = ""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        return text

    async def extract_text_from_docx(self, file_path: str) -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._sync_extract_docx_text, file_path)

    def _sync_extract_docx_text(self, file_path: str) -> str:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    async def extract_text_from_pptx(self, file_path: str) -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._sync_extract_pptx_text, file_path)

    def _sync_extract_pptx_text(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def clean_text(self, text: str) -> str:
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'[^\w\s.,;:!?()-]', '', text)
        return text.strip()

    async def summarize_text(
    self,
    text: str,
    method: str = "transformers",
    max_length: int = 130,
    min_length: int = 30,
    sentences_count: int = 3,
    tone: str = "formal"  
) -> str:
        if not text.strip():
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="No text content found"
            )

        try:
            if method == "transformers":
                loop = asyncio.get_event_loop()
                result = await loop.run_in_executor(
                    None,
                    lambda: self.summarizer_pipeline(
                        text,
                        max_length=max_length,
                        min_length=min_length,
                        do_sample=False
                    )
                )
                raw_summary = result[0]["summary_text"]

                return self.apply_tone(raw_summary, tone)

            else:
                summary = await self._sumy_summarize(text, sentences_count)
                return self.apply_tone(summary, tone)

        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Summarization failed: {str(e)}"
            )

    def apply_tone(self, summary: str, tone: str) -> str:
        """Modify summary based on tone"""
        if tone == "bullet":
            # Convert into bullet points (split by sentence)
            sentences = [s.strip() for s in summary.split('.') if s.strip()]
            return "\n".join(f"• {s}." for s in sentences)
        elif tone == "casual":
            # Add a slightly relaxed style (basic example)
            return summary.replace("However,", "But").replace("Moreover,", "Also,")
        # Default formal
        return summary


    async def _sumy_summarize(self, text: str, sentences_count: int) -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None,
            self._sync_sumy_summarize,
            text,
            sentences_count
        )

    def _sync_sumy_summarize(self, text: str, sentences_count: int) -> str:
        from sumy.parsers.plaintext import PlaintextParser
        from sumy.nlp.tokenizers import Tokenizer
        from sumy.summarizers.lsa import LsaSummarizer

        parser = PlaintextParser.from_string(text, Tokenizer("english"))
        summarizer = LsaSummarizer()
        summary = summarizer(parser.document, sentences_count)
        return " ".join(str(sentence) for sentence in summary)

    @staticmethod
    def validate_file_type(content_type: str, filename: str):
        ext = Path(filename).suffix.lower()
        if content_type not in ALLOWED_FILE_TYPES:
            raise HTTPException(400, "Unsupported file type")
        if ALLOWED_FILE_TYPES[content_type] != ext:
            raise HTTPException(400, "File extension doesn't match content type")

# Singleton instance
summarization_service = SummarizationService()
